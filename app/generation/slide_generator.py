"""Slide generator — creates PowerPoint presentations from PDF content."""

import os
import json
import asyncio
from pathlib import Path
from datetime import datetime
from typing import Any

from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

import cohere

from app.ingestion.text_embedder import embed_chunks
from app.ingestion.qdrant_client import search_figures_collection

load_dotenv()

COHERE_API_KEY = os.getenv("COHERE_API_KEY", "")
MODEL_NAME = "command-r-plus-08-2024"
OUTPUTS_DIR = Path("data/outputs")

DEFAULT_NUM_SLIDES = 10
DEFAULT_BULLETS = 4
DEFAULT_DEPTH = "detailed"
DEFAULT_STYLE = "academic"

DEPTH_INSTRUCTIONS = {
    "summary": (
        "Write SHORT bullet points — one line each. "
        "Overview only, no deep explanation."
    ),
    "detailed": (
        "Write DETAILED bullet points — 1-2 sentences each. "
        "Include explanations and context."
    ),
    "exam": (
        "Write bullet points as KEY FACTS to memorize. "
        "Include definitions, formulas, and important terms. "
        "Phrase as facts a student should know for an exam."
    )
}

STYLE_INSTRUCTIONS = {
    "academic": "Use formal academic language with precise terminology.",
    "simple": "Use plain simple language. Avoid jargon. Easy to understand.",
    "visual_hints": (
        "After each bullet, add a [VISUAL: ...] hint suggesting "
        "a diagram or image that would illustrate this point."
    )
}


async def _find_figures_for_slides(
    slides_data: list[dict[str, Any]],
    source_pdf: str,
    similarity_threshold: float = 0.45
) -> list[str | None]:
    """Pre-compute best figure paths for all slides concurrently.

    Batch-embeds all slide titles in one BGE-M3 call and searches
    Qdrant for all slides in parallel.

    Args:
        slides_data: List of slide dicts with title and chapter.
        source_pdf: PDF filename for Qdrant filter.
        similarity_threshold: Minimum combined score to accept a figure.

    Returns:
        List of absolute image paths (or None for slides without matches).
    """
    if not slides_data:
        return []

    queries = [
        f"{s.get('chapter', '')}: {s.get('title', '')}".strip(": ")
        for s in slides_data
    ]

    try:
        query_embeddings = await embed_chunks(queries)
    except Exception as e:
        print(f"[slide_generator] Batch embed error: {e}")
        return [None] * len(slides_data)

    async def _search_one(emb: dict) -> list[dict]:
        return await search_figures_collection(
            dense_vector=emb["dense_vector"],
            source_pdf=source_pdf,
            limit=5
        )

    all_results = await asyncio.gather(*[
        _search_one(emb) for emb in query_embeddings
    ])

    figures_dir = Path("data/processed/figures")
    image_paths: list[str | None] = []

    for idx, results in enumerate(all_results):
        if not results:
            image_paths.append(None)
            continue

        best_combined = 0.0
        best_result = None

        for result in results:
            desc_score = result.get("score", 0)
            combined = desc_score * 0.85
            if combined > best_combined:
                best_combined = combined
                best_result = result

        if best_result and best_combined >= similarity_threshold:
            figure_filename = best_result.get("figure_filename", "")
            figure_path = figures_dir / figure_filename if figure_filename else None
            if figure_path and figure_path.exists():
                print(f"[slide_generator] Figure for slide {idx}: {figure_path.name} (score={best_combined:.3f})")
                image_paths.append(str(figure_path))
            else:
                image_paths.append(None)
        else:
            image_paths.append(None)

    return image_paths


def _find_matching_image_fallback(
    slide_data: dict[str, Any],
    visual_chunks: list[dict[str, Any]],
    slide_index: int = 0,
    total_slides: int = 10
) -> str | None:
    """Fallback image matching using page distribution.

    Used when ColPali service is unavailable.

    Args:
        slide_data: Slide dict with title and chapter.
        visual_chunks: List of chunk payloads with screenshots.
        slide_index: Position of this slide (0-indexed).
        total_slides: Total number of slides being generated.

    Returns:
        Absolute path to PNG or None.
    """
    if not visual_chunks:
        return None

    valid = [
        c for c in visual_chunks
        if c.get("screenshot_path")
        and Path(c["screenshot_path"]).exists()
    ]
    if not valid:
        return None

    valid.sort(key=lambda c: c.get("page_num", 0))
    total_pages = len(valid)
    page_index = int((slide_index / max(total_slides, 1)) * total_pages)
    page_index = min(page_index, total_pages - 1)
    return valid[page_index]["screenshot_path"]


_cohere_client: cohere.AsyncClientV2 | None = None

def _get_cohere_client() -> cohere.AsyncClientV2:
    """Get or create async Cohere client (reused across calls)."""
    global _cohere_client
    if _cohere_client is None:
        _cohere_client = cohere.AsyncClientV2(api_key=COHERE_API_KEY)
    return _cohere_client


async def _extract_topics(
    chunks: list[dict[str, Any]],
    num_topics: int,
    topic_focus: str
) -> list[dict[str, Any]]:
    """Extract slide topics from chunks using Cohere.

    Args:
        chunks: List of chunk dicts with text and metadata.
        num_topics: Number of topics to extract.
        topic_focus: "all" or specific chapter/topic name.

    Returns:
        List of dicts with keys: topic, chapter, key_concept
    """
    if topic_focus != "all":
        filtered = [
            c for c in chunks
            if topic_focus.lower() in c.get("chapter", "").lower() or
            topic_focus.lower() in c.get("chunk", "").lower()
        ]
        context_chunks = filtered[:20] if filtered else chunks[:20]
    else:
        total = len(chunks)
        if total <= 30:
            context_chunks = chunks[:30]
        else:
            # Sample evenly across entire document
            step = max(1, total // 30)
            context_chunks = [chunks[i] for i in range(0, total, step)][:30]
            # Always include first and last chunks for context
            if chunks[0] not in context_chunks:
                context_chunks = [chunks[0]] + context_chunks[:29]
            if chunks[-1] not in context_chunks:
                context_chunks = context_chunks[:29] + [chunks[-1]]

    context = "\n\n".join([
        f"[{c.get('chapter', 'General')} | Page {c.get('page_num', 0)}]\n{c.get('chunk', '')[:300]}"
        for c in context_chunks
    ])

    prompt = f"""
You are analyzing a textbook. Extract exactly {num_topics} slide topics
from this content.

Content:
{context}

Return ONLY a JSON array like this:
[
  {{"topic": "Newton's Second Law", "chapter": "Chapter 3", "key_concept": "F = ma"}},
  {{"topic": "Circular Motion", "chapter": "Chapter 4", "key_concept": "centripetal force"}}
]

Rules:
- Exactly {num_topics} topics
- Topics should cover the most important concepts
- Each topic should be specific enough for one slide
- Return ONLY the JSON array, no other text
"""

    client = _get_cohere_client()
    response = await client.chat(
        model=MODEL_NAME,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=1000,
        temperature=0.3
    )
    raw = response.message.content[0].text

    try:
        raw = raw.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        topics = json.loads(raw.strip())
        return topics[:num_topics]
    except Exception as e:
        print(f"[slide_generator] Topic parse error: {e}")
        return [
            {"topic": f"Topic {i+1}", "chapter": "General", "key_concept": ""}
            for i in range(num_topics)
        ]


async def _generate_slide_content(
    topic: dict[str, Any],
    chunks: list[dict[str, Any]],
    bullets_per_slide: int,
    depth: str,
    style: str
) -> dict[str, Any]:
    """Generate content for a single slide using Cohere.

    Args:
        topic: Topic dict with topic, chapter, key_concept
        chunks: All chunks for context
        bullets_per_slide: Number of bullet points
        depth: "summary" | "detailed" | "exam"
        style: "academic" | "simple" | "visual_hints"

    Returns:
        Dict with keys: title, bullets, chapter
    """
    # Search for relevant chunks across entire document
    topic_words = set(topic.get("topic", "").lower().split())
    chapter_lower = topic.get("chapter", "").lower()

    scored = []
    for c in chunks:
        chunk_lower = c.get("chunk", "").lower()
        chapter_match = chapter_lower and chapter_lower in c.get("chapter", "").lower()
        word_matches = sum(1 for w in topic_words if w in chunk_lower and len(w) > 3)
        score = word_matches + (3 if chapter_match else 0)
        if score > 0:
            scored.append((score, c))

    scored.sort(key=lambda x: x[0], reverse=True)
    relevant = [c for _, c in scored[:5]]

    # Fallback: use evenly sampled chunks if no matches
    if not relevant:
        total = len(chunks)
        step = max(1, total // 5)
        relevant = [chunks[i] for i in range(0, total, step)][:5]

    context = "\n\n".join([
        c.get("chunk", "")[:400] for c in relevant
    ]) if relevant else "No specific context available."

    depth_inst = DEPTH_INSTRUCTIONS.get(depth, DEPTH_INSTRUCTIONS["detailed"])
    style_inst = STYLE_INSTRUCTIONS.get(style, STYLE_INSTRUCTIONS["academic"])

    prompt = f"""
Create content for a presentation slide about: {topic.get('topic', '')}
Chapter: {topic.get('chapter', '')}
Key concept: {topic.get('key_concept', '')}

Context from textbook:
{context}

{depth_inst}
{style_inst}

Return ONLY a JSON object like this:
{{
  "title": "Newton's Second Law",
  "bullets": [
    "Force equals mass times acceleration (F = ma)",
    "Doubling the force doubles the acceleration",
    "Doubling the mass halves the acceleration",
    "Direction of acceleration matches direction of net force"
  ]
}}

Rules:
- Exactly {bullets_per_slide} bullet points
- Title must be concise (max 8 words)
- Base content ONLY on the provided context
- Return ONLY the JSON object, no other text
"""

    client = _get_cohere_client()
    response = await client.chat(
        model=MODEL_NAME,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=600,
        temperature=0.4
    )
    raw = response.message.content[0].text

    try:
        raw = raw.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        content = json.loads(raw.strip())
        return {
            "title": content.get("title", topic.get("topic", "")),
            "bullets": content.get("bullets", [])[:bullets_per_slide],
            "chapter": topic.get("chapter", "")
        }
    except Exception as e:
        print(f"[slide_generator] Slide parse error: {e}")
        return {
            "title": topic.get("topic", "Slide"),
            "bullets": [f"Key point {i+1}" for i in range(bullets_per_slide)],
            "chapter": topic.get("chapter", "")
        }


async def _build_pptx(
    slides_data: list[dict[str, Any]],
    source_pdf: str,
    include_title_slide: bool,
    include_summary_slide: bool,
    slide_image_paths: list[str | None] | None = None
) -> str:
    """Build a .pptx file from slide content data.

    Args:
        slides_data: List of dicts with title and bullets
        source_pdf: Original PDF filename for title slide
        include_title_slide: Whether to add a title slide
        include_summary_slide: Whether to add a summary slide
        slide_image_paths: Pre-computed image paths (one per slide).
                           If None, no images are added.

    Returns:
        Absolute path to saved .pptx file as string
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
    TITLE_COLOR = RGBColor(0xE9, 0x4F, 0x37)
    TEXT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT_COLOR = RGBColor(0x39, 0x3E, 0x46)
    BULLET_COLOR = RGBColor(0xE0, 0xE0, 0xE0)

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    if include_title_slide:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        set_bg(slide, BG_COLOR)

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.33), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = Path(source_pdf).stem.replace("_", " ").title()
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER

        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2), Inches(11.33), Inches(0.8)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Generated by Online Teacher • {datetime.now().strftime('%B %Y')}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_COLOR
        p2.alignment = PP_ALIGN.CENTER

    for slide_index, slide_data in enumerate(slides_data):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        set_bg(slide, BG_COLOR)

        image_path = slide_image_paths[slide_index] if slide_image_paths else None
        has_image = image_path is not None

        if slide_data.get("chapter"):
            ch_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), Inches(12), Inches(0.4)
            )
            ch_tf = ch_box.text_frame
            ch_p = ch_tf.paragraphs[0]
            ch_p.text = slide_data["chapter"]
            ch_p.font.size = Pt(14)
            ch_p.font.color.rgb = ACCENT_COLOR

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.7), Inches(12.33), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR

        line = slide.shapes.add_shape(
            1,
            Inches(0.5), Inches(1.85),
            Inches(12.33), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = TITLE_COLOR
        line.line.fill.background()

        if has_image:
            # Two-column layout: bullets on left, image on right
            try:
                img = PILImage.open(image_path)
                img_width, img_height = img.size
                aspect = img_height / img_width

                max_w = Inches(5.6)
                max_h = Inches(5.0)
                img_w = max_w
                img_h = Emu(int(max_w * aspect))
                if img_h > max_h:
                    img_h = max_h
                    img_w = Emu(int(max_h / aspect))

                img_top = Inches(1.5) + (max_h - img_h) // 2

                slide.shapes.add_picture(
                    image_path,
                    left=Inches(7.2),
                    top=img_top,
                    width=img_w,
                    height=img_h
                )
                print(f"[slide_generator] Added image to slide: {slide_data['title']}")
            except Exception as e:
                print(f"[slide_generator] Could not add image: {e}")
                has_image = False

        # Determine content box width based on whether we have an image
        content_width = Inches(6.3) if has_image else Inches(11.93)

        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.1), content_width, Inches(4.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(slide_data.get("bullets", [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = BULLET_COLOR
            p.space_after = Pt(12)

    if include_summary_slide:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        set_bg(slide, BG_COLOR)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.7), Inches(12.33), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Takeaways"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR

        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.1), Inches(11.93), Inches(4.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, slide_data in enumerate(slides_data[:8]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {slide_data['title']}"
            p.font.size = Pt(20)
            p.font.color.rgb = BULLET_COLOR
            p.space_after = Pt(8)

    OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
    filename = f"{Path(source_pdf).stem}_slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = OUTPUTS_DIR / filename
    prs.save(str(output_path))
    print(f"[slide_generator] Saved: {output_path}")
    return str(output_path.resolve())


async def generate_slides(
    chunks: list[dict[str, Any]],
    source_pdf: str,
    num_slides: int = DEFAULT_NUM_SLIDES,
    bullets_per_slide: int = DEFAULT_BULLETS,
    depth: str = DEFAULT_DEPTH,
    topic_focus: str = "all",
    style: str = DEFAULT_STYLE,
    include_title_slide: bool = True,
    include_summary_slide: bool = True
) -> dict[str, Any]:
    """Generate a PowerPoint presentation from ingested PDF chunks.

    Args:
        chunks: All chunks from Qdrant for the source PDF.
        source_pdf: Original PDF filename.
        num_slides: Number of content slides to generate.
        bullets_per_slide: Bullet points per slide (3-5).
        depth: "summary" | "detailed" | "exam"
        topic_focus: "all" | chapter name | topic name
        style: "academic" | "simple" | "visual_hints"
        include_title_slide: Add title slide at start.
        include_summary_slide: Add summary slide at end.

    Returns:
        Dict with keys:
          - file_path: str (absolute path to .pptx)
          - num_slides: int (actual slides generated)
          - topics: list[str] (slide titles)
          - status: str ("success" or "error")
          - error: str | None
    """
    try:
        print(f"[slide_generator] Generating {num_slides} slides for {source_pdf}")

        topics = await _extract_topics(chunks, num_slides, topic_focus)
        print(f"[slide_generator] Extracted {len(topics)} topics")

            # Process in batches of 5 to avoid overwhelming Cohere
        async def _gather_in_batches(tasks, batch_size=5):
            results = []
            for i in range(0, len(tasks), batch_size):
                batch = tasks[i:i + batch_size]
                batch_results = await asyncio.gather(*batch)
                results.extend(batch_results)
            return results

        tasks = [
            _generate_slide_content(topic, chunks, bullets_per_slide, depth, style)
            for topic in topics
        ]
        slides_data = await _gather_in_batches(tasks, batch_size=5)                         

        print(f"[slide_generator] Generated content for {len(slides_data)} slides")

        # Pre-compute all figure paths in parallel before building slides
        slide_image_paths = await _find_figures_for_slides(
            slides_data, source_pdf
        )
        print(f"[slide_generator] Found figures for {sum(1 for p in slide_image_paths if p)}/{len(slides_data)} slides")

        file_path = await _build_pptx(
            list(slides_data),
            source_pdf,
            include_title_slide,
            include_summary_slide,
            slide_image_paths=slide_image_paths
        )

        return {
            "file_path": file_path,
            "num_slides": len(slides_data),
            "topics": [s["title"] for s in slides_data],
            "status": "success",
            "error": None
        }

    except Exception as e:
        print(f"[slide_generator] Error: {e}")
        return {
            "file_path": "",
            "num_slides": 0,
            "topics": [],
            "status": "error",
            "error": str(e)
        }